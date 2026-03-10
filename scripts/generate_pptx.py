"""
generate_pptx.py
Genera la presentacion "Orquestacion de agentes de IA" en formato .pptx
usando el contenido de plans/orquestacion-agentes-ia-deck.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os, re

# ---------------------------------------------------------------------------
# PALETA DE COLORES
# ---------------------------------------------------------------------------
C_DARK      = RGBColor(0x12, 0x16, 0x3E)   # azul noche
C_ACCENT    = RGBColor(0x7C, 0x3A, 0xED)   # violeta Copilot
C_ACCENT2   = RGBColor(0x06, 0xB6, 0xD4)   # cyan
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG  = RGBColor(0xF8, 0xF9, 0xFB)
C_BODY_TEXT = RGBColor(0x1F, 0x29, 0x37)
C_GRAY      = RGBColor(0x6B, 0x72, 0x80)
C_ACCENT_L  = RGBColor(0xED, 0xE9, 0xFE)   # violeta claro (fondo bullets especiales)

# ---------------------------------------------------------------------------
# FUENTES
# ---------------------------------------------------------------------------
FONT_TITLE  = "Segoe UI"
FONT_BODY   = "Segoe UI"

# ---------------------------------------------------------------------------
# TAMANOS DE SLIDE (widescreen 16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                font_size=18, bold=False, color=C_BODY_TEXT,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# TIPOS DE SLIDE
# ---------------------------------------------------------------------------

def make_portada(prs, title_text, subtitle_text):
    """Slide de portada: fondo oscuro, titulo grande, linea de acento."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, C_DARK)

    # Banda decorativa superior
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.55), C_ACCENT)

    # Linea cyan en la seccion de titulo
    add_rect(slide, Inches(0.8), Inches(3.3), Inches(0.06), Inches(1.4), C_ACCENT2)

    # Titulo principal
    add_textbox(slide,
                left=Inches(1.0), top=Inches(1.6),
                width=Inches(11.0), height=Inches(1.4),
                text=title_text,
                font_name=FONT_TITLE, font_size=36, bold=True,
                color=C_WHITE, align=PP_ALIGN.LEFT)

    # Subtitulo
    add_textbox(slide,
                left=Inches(1.0), top=Inches(3.2),
                width=Inches(10.5), height=Inches(0.8),
                text=subtitle_text,
                font_name=FONT_BODY, font_size=18, bold=False,
                color=C_ACCENT2, align=PP_ALIGN.LEFT)

    # Etiqueta inferior
    add_textbox(slide,
                left=Inches(1.0), top=Inches(6.6),
                width=Inches(8), height=Inches(0.5),
                text="Accenture · AI Engineering Guild · Marzo 2026",
                font_name=FONT_BODY, font_size=12,
                color=RGBColor(0xA0, 0xAE, 0xC0), align=PP_ALIGN.LEFT)

    # Banda decorativa inferior
    add_rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), C_ACCENT2)
    return slide


def make_agenda(prs, blocks):
    """Slide de agenda con bloques numerados."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_LIGHT_BG)

    # Header
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), C_DARK)
    add_textbox(slide,
                Inches(0.5), Inches(0.25), Inches(10), Inches(0.7),
                "AGENDA", FONT_TITLE, 28, True, C_WHITE, PP_ALIGN.LEFT)
    add_rect(slide, Inches(0.5), Inches(1.15), Inches(2), Inches(0.05), C_ACCENT)

    # Bloques numerados
    x_left = Inches(0.7)
    y_start = Inches(1.6)
    row_h   = Inches(0.85)
    num_w   = Inches(0.55)
    txt_w   = Inches(10.8)

    for i, bloque in enumerate(blocks):
        y = y_start + i * row_h
        # circulo/numero
        add_rect(slide, x_left, y + Inches(0.08), num_w, num_w, C_ACCENT)
        add_textbox(slide, x_left, y + Inches(0.05), num_w, num_w,
                    str(i + 1), FONT_TITLE, 14, True, C_WHITE, PP_ALIGN.CENTER)
        # texto del bloque
        add_textbox(slide, x_left + Inches(0.7), y, txt_w, row_h,
                    bloque, FONT_BODY, 16, False, C_BODY_TEXT, PP_ALIGN.LEFT)

    # Franja inferior accent
    add_rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), C_ACCENT)
    return slide


def make_section_divider(prs, block_label, block_title):
    """Slide divisor de seccion (fondo oscuro con etiqueta del bloque)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK)

    # Linea decorativa
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.35), C_ACCENT)
    add_rect(slide, Inches(1.2), Inches(3.0), Inches(10.5), Inches(0.04), C_ACCENT2)

    # Etiqueta del bloque
    add_textbox(slide,
                Inches(1.2), Inches(2.0), Inches(10), Inches(0.8),
                block_label.upper(),
                FONT_TITLE, 16, True, C_ACCENT2, PP_ALIGN.LEFT)

    # Titulo del bloque
    add_textbox(slide,
                Inches(1.2), Inches(2.7), Inches(10.5), Inches(2.0),
                block_title,
                FONT_TITLE, 34, True, C_WHITE, PP_ALIGN.LEFT)

    add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), C_ACCENT2)
    return slide


def make_content_slide(prs, title, bullets, tiempo=None, slide_num=None):
    """Slide de contenido: header oscuro + bullets jerarquicos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_LIGHT_BG)

    # Header
    header_h = Inches(1.3)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, header_h, C_DARK)

    # Titulo en header
    add_textbox(slide,
                Inches(0.5), Inches(0.2), Inches(10.5), Inches(0.85),
                title, FONT_TITLE, 22, True, C_WHITE, PP_ALIGN.LEFT)

    # Badge de tiempo
    if tiempo:
        add_rect(slide,
                 Inches(11.5), Inches(0.3), Inches(1.5), Inches(0.55),
                 C_ACCENT)
        add_textbox(slide,
                    Inches(11.5), Inches(0.3), Inches(1.5), Inches(0.55),
                    tiempo, FONT_BODY, 11, False, C_WHITE, PP_ALIGN.CENTER)

    # Linea decorativa bajo header
    add_rect(slide, Inches(0), header_h, SLIDE_W, Inches(0.04), C_ACCENT)

    # === Bullets ===
    y_cur   = Inches(1.5)
    margin_l = Inches(0.55)
    txt_w    = Inches(12.0)
    max_y    = Inches(7.0)

    for bullet in bullets:
        if y_cur >= max_y:
            break

        level = bullet.get("level", 0)
        text  = bullet.get("text", "")
        is_key = bullet.get("key", False)

        indent = Inches(0.4 * level)
        fsize  = 15 - (level * 1)
        fsize  = max(fsize, 11)
        bh     = Inches(0.38) if level == 0 else Inches(0.32)

        # marcador de bullet
        if level == 0:
            marker_color = C_ACCENT
            marker_w = Inches(0.12)
            marker_h = Inches(0.12)
            add_rect(slide,
                     margin_l + indent,
                     y_cur + Inches(0.12),
                     marker_w, marker_h, marker_color)
            text_left = margin_l + indent + Inches(0.22)
        else:
            marker_color = C_ACCENT2
            marker_w = Inches(0.08)
            marker_h = Inches(0.08)
            add_rect(slide,
                     margin_l + indent + Inches(0.05),
                     y_cur + Inches(0.11),
                     marker_w, marker_h, marker_color)
            text_left = margin_l + indent + Inches(0.22)

        available_w = txt_w - indent - Inches(0.22)
        # fondo tenue para bullets de nivel 0 "destacados"
        if is_key and level == 0:
            add_rect(slide,
                     margin_l + indent, y_cur - Inches(0.03),
                     available_w + Inches(0.1), bh + Inches(0.06),
                     C_ACCENT_L)

        add_textbox(slide,
                    text_left, y_cur,
                    available_w, bh + Inches(0.1),
                    text, FONT_BODY, fsize,
                    bold=(level == 0),
                    color=C_BODY_TEXT if level > 0 else C_BODY_TEXT,
                    align=PP_ALIGN.LEFT)

        y_cur += bh + Inches(0.05)

    # Numero de slide
    if slide_num:
        add_textbox(slide,
                    Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3),
                    str(slide_num), FONT_BODY, 10, False, C_GRAY, PP_ALIGN.RIGHT)

    # Franja inferior
    add_rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), C_ACCENT2)
    return slide


def make_cierre(prs, bullets):
    """Slide de cierre con fondo oscuro y consejos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.35), C_ACCENT)

    add_textbox(slide,
                Inches(0.8), Inches(0.55), Inches(11), Inches(0.9),
                "CIERRE · DEMO EN VIVO · Q&A",
                FONT_TITLE, 26, True, C_WHITE, PP_ALIGN.LEFT)

    add_rect(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.04), C_ACCENT2)

    y_cur = Inches(1.6)
    for bullet in bullets:
        level = bullet.get("level", 0)
        text  = bullet.get("text", "")
        indent = Inches(0.4 * level)
        fsize = 15 - level
        fsize = max(fsize, 11)
        bh = Inches(0.4)

        color = C_ACCENT2 if level == 0 else RGBColor(0xD1, 0xD5, 0xDB)
        add_textbox(slide,
                    Inches(1.0) + indent, y_cur,
                    Inches(11.5) - indent, bh,
                    ("▸  " if level == 0 else "   — ") + text,
                    FONT_BODY, fsize, bold=(level == 0),
                    color=color, align=PP_ALIGN.LEFT)
        y_cur += bh + Inches(0.04)

    add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), C_ACCENT)
    return slide


# ---------------------------------------------------------------------------
# PARSER del MARKDOWN
# ---------------------------------------------------------------------------

def parse_bullets(lines):
    """Convierte lineas de bullets Markdown a lista de dicts {level, text}."""
    bullets = []
    for raw in lines:
        stripped = raw.rstrip()
        if not stripped:
            continue
        # contar espacios de indentacion (2 o 4 => level 1, 4+ => level 2)
        leading = len(raw) - len(raw.lstrip())
        if leading == 0:
            level = 0
        elif leading <= 4:
            level = 1
        else:
            level = 2
        # quitar el "- " o "  - "
        text = re.sub(r"^\s*[-*]\s*", "", stripped)
        # quitar backticks de codigo inline
        text = text.replace("`", "")
        bullets.append({"level": level, "text": text})
    return bullets


def parse_deck_md(filepath):
    """Parsea el Markdown y devuelve lista de slides {num, title, tiempo, bullets}."""
    with open(filepath, encoding="utf-8") as f:
        content = f.read()

    # Quitar bloque markdown wrappers si los hay
    content = re.sub(r"^```markdown\s*\n", "", content)
    content = re.sub(r"\n```\s*$", "", content)

    slides_raw = re.split(r"\n## Slide \d+", content)
    slides = []

    for i, block in enumerate(slides_raw[1:], start=1):  # saltar header global
        lines = block.split("\n")
        # Primera linea: " — Titulo"
        title_line = lines[0].strip() if lines else ""
        title = re.sub(r"^\s*—\s*", "", title_line).strip()

        bullet_lines = []
        tiempo = None
        for line in lines[1:]:
            if line.strip().startswith("- Tiempo:"):
                tiempo = line.strip().replace("- Tiempo:", "").strip()
            elif re.match(r"^\s*[-*]", line):
                bullet_lines.append(line)

        bullets = parse_bullets(bullet_lines)
        slides.append({
            "num": i,
            "title": title,
            "tiempo": tiempo,
            "bullets": bullets,
        })

    return slides


# ---------------------------------------------------------------------------
# DETECTAR BLOQUES para divisores de seccion
# ---------------------------------------------------------------------------

BLOQUE_MAP = {
    "Bloque 1": ("Bloque 1", "El agente monolítico"),
    "Bloque 2": ("Bloque 2", "Qué es orquestación · Roles · Modelos"),
    "Bloque 3": ("Bloque 3", "Handoffs · IDEs · Configuración práctica"),
    "Bloque 4": ("Bloque 4", "Puntos de control (Gates) · Flujo end-to-end"),
    "Bloque 5": ("Bloque 5", "Buenas prácticas · Anti-patrones · Cierre"),
}

def get_block_key(title):
    for key in BLOQUE_MAP:
        if key in title:
            return key
    return None


# ---------------------------------------------------------------------------
# MAIN: CONSTRUIR EL PPTX
# ---------------------------------------------------------------------------

def build_pptx(md_path, out_path):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides_data = parse_deck_md(md_path)
    seen_blocks = set()

    for s in slides_data:
        num    = s["num"]
        title  = s["title"]
        tiempo = s["tiempo"]
        bullets = s["bullets"]

        # --- PORTADA (Slide 1) ---
        if num == 1:
            make_portada(prs,
                         "Orquestación de agentes de IA\nen desarrollo de software",
                         "Roles · Handoffs · Gates · Token Economics")
            continue

        # --- AGENDA (Slide 2) ---
        if num == 2:
            agenda_items = [b["text"] for b in bullets if b["level"] == 0]
            make_agenda(prs, agenda_items)
            continue

        # --- DIVISOR DE SECCION (cuando hay un nuevo bloque) ---
        block_key = get_block_key(title)
        if block_key and block_key not in seen_blocks:
            seen_blocks.add(block_key)
            lbl, desc = BLOQUE_MAP[block_key]
            make_section_divider(prs, lbl, desc)

        # --- CIERRE (Slide 11) ---
        if num == len(slides_data):
            make_cierre(prs, bullets)
            continue

        # --- SLIDE DE CONTENIDO NORMAL ---
        make_content_slide(prs, title, bullets, tiempo, slide_num=num)

    prs.save(out_path)
    print(f"\n✅  Presentación generada: {out_path}")
    print(f"   Slides generadas: {len(prs.slides)}")


# ---------------------------------------------------------------------------
# ENTRY POINT
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    md   = os.path.join(base, "plans", "orquestacion-agentes-ia-deck.md")
    out_dir = os.path.join(base, "out")
    os.makedirs(out_dir, exist_ok=True)
    out  = os.path.join(out_dir, "orquestacion-agentes-ia.pptx")
    build_pptx(md, out)
